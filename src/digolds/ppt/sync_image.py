from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def _parse(args):
    ppt_file = args.get('--ppt-file')
    slide_no = args.get('--slide-no')
    image_name = args.get('--image-name')
    new_image_file = args.get('--new-image-file')
    return [ppt_file, slide_no, image_name, new_image_file]

def _update_image(ppt_file, slide_no, image_name, new_image_file):
    prs = Presentation(ppt_file)
    slide = prs.slides[slide_no]
    shapes = slide.shapes
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and os.path.basename(slide.part.rels.related_parts[s._pic.blip_rId].partname) == image_name:
            """
            left = s.left
            top = s.top
            width = s.width
            height = s.height
            new_shape = shapes.add_picture(new_image_file,left,top, width, height)
            old_pic = s._element
            new_pic = new_shape._element
            old_pic.addnext(new_pic)
            old_pic.getparent().remove(old_pic)
            """
            # noinspection PyProtectedMember
            with open(new_image_file, "rb") as file_obj:
                r_img_blob = file_obj.read()
            img_rid = s._pic.blip_rId
            img_part = slide.part.related_parts[img_rid]
            img_part._blob = r_img_blob
            break
    prs.save(ppt_file)

name = 'update-image'

def operator(df, args):
    return _update_image(*_parse(args))

if __name__ == "__main__":
    import pathlib

    args = {
        '--ppt-file':str(pathlib.Path.cwd() /'ppt/tests/sample.pptx'),
        '--slide-no': 0,
        '--image-name' : 'Icon_Bird_512x512.png',
        '--new-image-file' : str(pathlib.Path.cwd() / 'ppt/tests/index.jpg')
    }

    operator(None, args)

# https://stackoverflow.com/questions/46463861/get-image-file-names-with-python-pptx
# https://towardsdatascience.com/creating-presentations-with-python-3f5737824f61
# https://python-pptx.readthedocs.io/en/latest/api/presentation.html#presentation-function