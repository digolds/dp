from pptx import Presentation

def _parse(args):
    ppt_file = args.get('--ppt-file')
    slide_no = args.get('--slide-no')
    shape_id = args.get('--shape-id')
    new_text = args.get('--new-text')
    return [ppt_file, slide_no, shape_id, new_text]

def _update_text(ppt_file, slide_no, shape_id, new_text):
    prs = Presentation(ppt_file)
    slide = prs.slides[slide_no]
    shapes = slide.shapes
    for s in shapes:
        if s.has_text_frame and s.shape_id == shape_id:
            s.text = new_text
            break
    prs.save(ppt_file)

name = 'update-text'

def operator(df, args):
    return _update_text(*_parse(args))

if __name__ == "__main__":
    import pathlib

    args = {
        '--ppt-file': r"D:\dr\src\python-lab\sample.pptx",
        '--slide-no': 1,
        '--shape-id' : 4,
        '--new-text' : 'I love digolds'
    }

    operator(None, args)

# https://stackoverflow.com/questions/46463861/get-image-file-names-with-python-pptx
# https://towardsdatascience.com/creating-presentations-with-python-3f5737824f61
# https://python-pptx.readthedocs.io/en/latest/api/presentation.html#presentation-function