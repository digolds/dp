from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

def _parse(args):
    ppt_file = args.get('--ppt-file')
    slide_no = args.get('--slide-no')
    image_name = args.get('--image-name')
    new_image_file = args.get('--new-image-file')
    return [ppt_file, slide_no, image_name, new_image_file]

def _create_ppt():
    path='ppt/tests/Icon_Bird_512x512.png'
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left=Inches(1)
    top=Inches(0.5)
    slide.shapes.add_picture(path,left,top)
    prs.save('ppt/tests/sample.pptx')

def _update_image(ppt_file, slide_no, image_name, new_image_file):
    prs = Presentation(ppt_file)
    for s in prs.slides[slide_no].shapes:
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and s._pic.nvPicPr.cNvPr.get('descr') == image_name:
            left = s.left
            top = s.top
            width = s.width
            height = s.height
            new_shape = prs.slides[slide_no].shapes.add_picture(new_image_file,left,top, width, height)
            old_pic = s._element
            new_pic = new_shape._element
            old_pic.addnext(new_pic)
            old_pic.getparent().remove(old_pic)
            break
    prs.save(ppt_file)

name = 'update-pptx'

def operator(df, args):
    return _update_image(*_parse(args))

if __name__ == "__main__":
    args = {
        '--ppt-file':r'ppt/tests/sample.pptx',
        '--slide-no': 0,
        '--image-name' : 'Icon_Bird_512x512.png',
        '--new-image-file' : 'ppt/tests/index.jpg'
    }
    _create_ppt()
    operator(None, args)

# https://stackoverflow.com/questions/46463861/get-image-file-names-with-python-pptx
# https://towardsdatascience.com/creating-presentations-with-python-3f5737824f61
# https://python-pptx.readthedocs.io/en/latest/api/presentation.html#presentation-function